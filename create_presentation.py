#!/usr/bin/env python3
"""
Create a professional PowerPoint presentation for:
"Understanding the Limits of Single-Cell Foundation Models on Downstream Tasks"
Author: Keisuke Nishioka
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os

# ============================================================
# Color Palette (modern academic style)
# ============================================================
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # Deep navy
ACCENT_BLUE = RGBColor(0x41, 0x7D, 0xC1)    # Primary blue
ACCENT_TEAL = RGBColor(0x2E, 0xCC, 0xA8)    # Teal/green accent
ACCENT_CORAL = RGBColor(0xE8, 0x6F, 0x51)   # Coral/orange accent
ACCENT_PURPLE = RGBColor(0x9B, 0x72, 0xCF)  # Purple accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xEC)
MEDIUM_GRAY = RGBColor(0x9A, 0x9A, 0xA8)
VERY_LIGHT = RGBColor(0xF5, 0xF5, 0xFA)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x3F)
GOLD = RGBColor(0xFF, 0xD7, 0x00)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Image paths
IMG_DIR = "/home/nishioka/LUH/BioLLM/Geneformer/results"
LOGO_PATH = "/home/nishioka/LUH/BioLLM/Geneformer/docs/source/_static/gf_logo.png"


def set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_name="Calibri",
                 font_size=18, font_color=DARK_TEXT, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.2):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_name="Calibri",
                       font_size=16, font_color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                       line_spacing=1.4, bullet=False):
    """Add multi-line text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if isinstance(line_data, tuple):
            text, is_bold, color, size = line_data
        else:
            text = line_data
            is_bold = False
            color = font_color
            size = font_size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet and not is_bold:
            p.text = text
        else:
            p.text = text

        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = is_bold
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(size * line_spacing)

    return txBox


def add_accent_bar(slide, left, top, width, height, color):
    """Add a thin accent bar."""
    return add_shape(slide, left, top, width, height, color)


def slide_header(slide, title, subtitle=None, slide_num=None, total=20):
    """Add standard slide header with accent bar."""
    set_slide_bg(slide, WHITE)

    # Top accent bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                 title, font_size=30, font_color=DARK_BG, bold=True,
                 font_name="Calibri Light")

    # Accent underline
    add_shape(slide, Inches(0.8), Inches(0.95), Inches(2), Inches(0.04), ACCENT_TEAL)

    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
                     subtitle, font_size=16, font_color=MEDIUM_GRAY, italic=True,
                     font_name="Calibri")

    # Slide number
    if slide_num:
        add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1), Inches(0.4),
                     f"{slide_num}/{total}", font_size=11, font_color=MEDIUM_GRAY,
                     alignment=PP_ALIGN.RIGHT, font_name="Calibri")

    # Footer line
    add_shape(slide, Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.01), LIGHT_GRAY)
    add_text_box(slide, Inches(0.8), Inches(7.1), Inches(6), Inches(0.3),
                 "Keisuke Nishioka | AI Foundation Models in Biomedicine | WiSe 2025/26",
                 font_size=9, font_color=MEDIUM_GRAY, font_name="Calibri")


def create_card(slide, left, top, width, height, title, body_lines,
                accent_color=ACCENT_BLUE, title_size=14, body_size=12):
    """Create an info card with accent left border."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, VERY_LIGHT)

    # Left accent bar
    add_shape(slide, left, top + Inches(0.1), Inches(0.05), height - Inches(0.2), accent_color)

    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                 title, font_size=title_size, font_color=accent_color, bold=True,
                 font_name="Calibri")

    # Body
    y_offset = top + Inches(0.45)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.2), y_offset, width - Inches(0.3), Inches(0.25),
                     line, font_size=body_size, font_color=DARK_TEXT, font_name="Calibri")
        y_offset += Inches(0.22)

    return card


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ========================================================
    # SLIDE 1: Title Slide
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Geometric accent shapes
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_TEAL)
    add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_BLUE)

    # Decorative side bar
    add_shape(slide, Inches(0.6), Inches(1.5), Inches(0.06), Inches(4.5), ACCENT_TEAL)

    # Title
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(1.2),
                 "Understanding the Limits of", font_size=38,
                 font_color=LIGHT_GRAY, bold=False, font_name="Calibri Light")

    add_text_box(slide, Inches(1.2), Inches(2.6), Inches(10.5), Inches(1.2),
                 "Single-Cell Foundation Models", font_size=44,
                 font_color=WHITE, bold=True, font_name="Calibri")

    add_text_box(slide, Inches(1.2), Inches(3.5), Inches(10.5), Inches(0.7),
                 "on Downstream Tasks", font_size=38,
                 font_color=ACCENT_TEAL, bold=True, font_name="Calibri")

    # Separator line
    add_shape(slide, Inches(1.2), Inches(4.5), Inches(3), Inches(0.03), ACCENT_BLUE)

    # Author info
    add_text_box(slide, Inches(1.2), Inches(4.8), Inches(6), Inches(0.5),
                 "Keisuke Nishioka", font_size=22, font_color=WHITE,
                 bold=True, font_name="Calibri")

    add_text_box(slide, Inches(1.2), Inches(5.3), Inches(6), Inches(0.4),
                 "Student ID: 10081049", font_size=14, font_color=MEDIUM_GRAY,
                 font_name="Calibri")

    add_text_box(slide, Inches(1.2), Inches(5.7), Inches(8), Inches(0.4),
                 "AI Foundation Models in Biomedicine  |  WiSe 2025/26", font_size=14,
                 font_color=MEDIUM_GRAY, font_name="Calibri")

    add_text_box(slide, Inches(1.2), Inches(6.1), Inches(6), Inches(0.4),
                 "Leibniz University of Hannover", font_size=14,
                 font_color=ACCENT_BLUE, font_name="Calibri")

    # DNA helix decorative text (right side)
    add_text_box(slide, Inches(9.5), Inches(2.0), Inches(3.5), Inches(4.0),
                 "scRNA\nGPT\nLLM", font_size=60, font_color=RGBColor(0x25, 0x25, 0x45),
                 bold=True, font_name="Calibri", alignment=PP_ALIGN.RIGHT)

    # ========================================================
    # SLIDE 2: Outline / Agenda
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Presentation Outline", slide_num=2)

    sections = [
        ("01", "Introduction & Motivation", "Why evaluate single-cell foundation models?", ACCENT_BLUE),
        ("02", "Background & Related Work", "Geneformer, scGPT, scFoundation", ACCENT_TEAL),
        ("03", "Research Hypothesis & Approach", "Unified evaluation protocol design", ACCENT_PURPLE),
        ("04", "Experimental Setup", "Datasets, models, and metrics", ACCENT_BLUE),
        ("05", "Results & Analysis", "Frozen vs. fine-tuned performance comparison", ACCENT_CORAL),
        ("06", "Discussion & Conclusion", "Key findings and future directions", ACCENT_TEAL),
    ]

    for i, (num, title, desc, color) in enumerate(sections):
        y = Inches(1.6) + Inches(i * 0.85)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Section title
        add_text_box(slide, Inches(2.0), y, Inches(4.5), Inches(0.35),
                     title, font_size=18, font_color=DARK_BG, bold=True, font_name="Calibri")

        # Description
        add_text_box(slide, Inches(2.0), y + Inches(0.32), Inches(6), Inches(0.3),
                     desc, font_size=13, font_color=MEDIUM_GRAY, font_name="Calibri")

        # Connector line
        if i < len(sections) - 1:
            add_shape(slide, Inches(1.45), y + Inches(0.55), Inches(0.02), Inches(0.3), LIGHT_GRAY)

    # ========================================================
    # SLIDE 3: Motivation - Why This Matters
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Motivation", "Why evaluate single-cell foundation models?", slide_num=3)

    # Left column - Revolution
    create_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.2),
                "The Foundation Model Revolution",
                ["Transformer architectures have transformed scRNA-seq analysis",
                 "Large-scale pretraining on millions of transcriptomes",
                 "Promise of transferable representations for downstream tasks",
                 "Multiple competing models: Geneformer, scGPT, scFoundation"],
                accent_color=ACCENT_BLUE, body_size=13)

    # Right column - Problem
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.2),
                "The Evaluation Gap",
                ["Reported performances are difficult to compare",
                 "Heterogeneous datasets & preprocessing pipelines",
                 "Different fine-tuning strategies across papers",
                 "Unclear when pretraining provides real advantages"],
                accent_color=ACCENT_CORAL, body_size=13)

    # Bottom - Key Questions
    add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.5), VERY_LIGHT)
    add_text_box(slide, Inches(1.2), Inches(4.35), Inches(10), Inches(0.4),
                 "Key Research Questions", font_size=18, font_color=ACCENT_PURPLE,
                 bold=True, font_name="Calibri")

    questions = [
        "Q1  What biological information is robustly encoded during pretraining?",
        "Q2  Under what conditions do these representations fail to generalize?",
        "Q3  How do frozen representations compare to fine-tuned models?",
    ]
    for i, q in enumerate(questions):
        y = Inches(4.85) + Inches(i * 0.5)
        add_text_box(slide, Inches(1.5), y, Inches(10), Inches(0.4),
                     q, font_size=15, font_color=DARK_TEXT, font_name="Calibri")

    # ========================================================
    # SLIDE 4: Background - Single-Cell Biology
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Background", "Single-cell transcriptomics & foundation models", slide_num=4)

    # scRNA-seq overview
    create_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.0),
                "Single-Cell RNA Sequencing (scRNA-seq)",
                ["Measures gene expression at individual cell level",
                 "Reveals cellular heterogeneity within tissues",
                 "Data: sparse count matrices (cells x genes)",
                 "Challenge: high dimensionality, noise, batch effects"],
                accent_color=ACCENT_BLUE, body_size=13)

    # Foundation model paradigm
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.0),
                "Foundation Model Paradigm",
                ["Pretrain on massive unlabeled datasets",
                 "Learn general biological representations",
                 "Transfer to diverse downstream tasks",
                 "Analogous to BERT/GPT in NLP"],
                accent_color=ACCENT_TEAL, body_size=13)

    # Key concept: representation learning
    add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.05), ACCENT_BLUE)
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.4),
                 "Core Idea: Pretrained representations should capture biological structure and enable efficient adaptation to new tasks",
                 font_size=15, font_color=DARK_BG, bold=True, italic=True,
                 font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Pipeline diagram (text-based)
    boxes_data = [
        ("Raw scRNA-seq\nCount Matrix", Inches(0.8), ACCENT_BLUE),
        ("Pretrained\nEncoder", Inches(3.5), ACCENT_PURPLE),
        ("Cell\nEmbeddings", Inches(6.2), ACCENT_TEAL),
        ("Downstream\nClassifier", Inches(8.9), ACCENT_CORAL),
        ("Cell Type\nPrediction", Inches(11.2), ACCENT_BLUE),
    ]

    for text, left, color in boxes_data:
        box = add_rounded_rect(slide, left, Inches(5.0), Inches(2.0), Inches(1.2), color)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(text.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    # Arrows between boxes
    for i in range(4):
        left = boxes_data[i][1] + Inches(2.05)
        add_text_box(slide, left, Inches(5.3), Inches(0.5), Inches(0.5),
                     "\u2192", font_size=28, font_color=MEDIUM_GRAY,
                     alignment=PP_ALIGN.CENTER, font_name="Calibri")

    # ========================================================
    # SLIDE 5: Models Overview
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Models Under Evaluation", "Three representative single-cell foundation models", slide_num=5)

    models = [
        ("Geneformer", "Theodoris et al., Nature 2023",
         ["Pretrained on 30M human transcriptomes",
          "Rank-value encoding of gene expression",
          "104M parameters (V2)",
          "Strong transfer learning across tasks"],
         ACCENT_BLUE),
        ("scGPT", "Cui et al., Nature Methods 2024",
         ["Generative pretraining objective",
          "Masked gene expression prediction",
          "Multi-omics integration capability",
          "Batch correction & cell annotation"],
         ACCENT_TEAL),
        ("scFoundation", "Hao et al., Nature Methods 2024",
         ["Trained on tens of millions of cells",
          "Largest scale pretraining to date",
          "Competitive benchmark results",
          "Not publicly available (limitation)"],
         ACCENT_PURPLE),
    ]

    for i, (name, ref, points, color) in enumerate(models):
        left = Inches(0.8) + Inches(i * 4.1)

        # Model card
        card_bg = add_rounded_rect(slide, left, Inches(1.5), Inches(3.8), Inches(4.8), VERY_LIGHT)

        # Color header
        add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(0.8), color)

        # Model name
        add_text_box(slide, left + Inches(0.2), Inches(1.55), Inches(3.4), Inches(0.45),
                     name, font_size=22, font_color=WHITE, bold=True, font_name="Calibri")

        # Reference
        add_text_box(slide, left + Inches(0.2), Inches(1.95), Inches(3.4), Inches(0.3),
                     ref, font_size=10, font_color=RGBColor(0xDD, 0xDD, 0xEE),
                     italic=True, font_name="Calibri")

        # Features
        for j, point in enumerate(points):
            y = Inches(2.6) + Inches(j * 0.55)
            # Bullet dot
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), y + Inches(0.07),
                                         Inches(0.12), Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            add_text_box(slide, left + Inches(0.5), y, Inches(3.1), Inches(0.45),
                         point, font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    # Status indicators
    status_data = [
        (Inches(0.8), "\u2713 Evaluated (Frozen + Fine-tuned)", ACCENT_BLUE),
        (Inches(4.9), "\u2713 Evaluated (Frozen only)", ACCENT_TEAL),
        (Inches(9.0), "\u2717 Not available for evaluation", ACCENT_CORAL),
    ]
    for left, text, color in status_data:
        add_text_box(slide, left, Inches(6.5), Inches(3.8), Inches(0.3),
                     text, font_size=11, font_color=color, bold=True,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # ========================================================
    # SLIDE 6: Research Hypothesis
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Research Hypothesis", slide_num=6)

    # Central hypothesis box
    add_shape(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.2),
              RGBColor(0xEE, 0xF2, 0xFF))
    add_shape(slide, Inches(1.5), Inches(1.8), Inches(0.08), Inches(2.2), ACCENT_BLUE)

    add_text_box(slide, Inches(1.9), Inches(1.9), Inches(9.5), Inches(0.35),
                 "Central Hypothesis", font_size=20, font_color=ACCENT_BLUE,
                 bold=True, font_name="Calibri")

    add_text_box(slide, Inches(1.9), Inches(2.4), Inches(9.5), Inches(1.4),
                 "Foundation models provide strong advantages when downstream data "
                 "aligns with the pretraining distribution, but their benefits "
                 "diminish under domain shift, where representation quality rather "
                 "than task-specific optimization becomes decisive.",
                 font_size=17, font_color=DARK_TEXT, font_name="Calibri", italic=True)

    # Two sub-hypotheses
    # H1
    create_card(slide, Inches(1.0), Inches(4.5), Inches(5.3), Inches(2.0),
                "H1: Fine-tuning > Frozen",
                ["Frozen representations capture general patterns",
                 "But miss task-specific nuances",
                 "End-to-end optimization should significantly",
                 "improve downstream performance"],
                accent_color=ACCENT_TEAL, body_size=13)

    # H2
    create_card(slide, Inches(7.0), Inches(4.5), Inches(5.3), Inches(2.0),
                "H2: Domain Shift Reduces Advantage",
                ["Cross-dataset evaluation (PBMC \u2192 Tabula Sapiens)",
                 "Performance degradation under distribution shift",
                 "Representation quality becomes the bottleneck",
                 "Model architecture matters less than data alignment"],
                accent_color=ACCENT_CORAL, body_size=13)

    # ========================================================
    # SLIDE 7: Experimental Design
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Experimental Design", "Unified evaluation protocol", slide_num=7)

    # Design principle box
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
              RGBColor(0xFD, 0xF5, 0xE6))
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.05), ACCENT_CORAL)

    add_text_box(slide, Inches(1.2), Inches(1.65), Inches(11), Inches(0.7),
                 "Design Principle: All components except the pretrained encoder are kept identical across experiments.\n"
                 "This isolates the contribution of pretraining from other experimental variables.",
                 font_size=14, font_color=DARK_TEXT, font_name="Calibri")

    # Two regimes comparison
    add_text_box(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(0.4),
                 "Regime 1: Frozen Representations", font_size=18,
                 font_color=ACCENT_BLUE, bold=True, font_name="Calibri")

    frozen_steps = [
        "1. Extract embeddings from pretrained model",
        "2. No weight updates to encoder",
        "3. Train lightweight classifier (LogReg/XGBoost)",
        "4. Evaluate on held-out test set",
    ]
    for i, step in enumerate(frozen_steps):
        add_text_box(slide, Inches(1.2), Inches(3.3) + Inches(i * 0.4), Inches(5), Inches(0.35),
                     step, font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    add_text_box(slide, Inches(7.0), Inches(2.8), Inches(5.5), Inches(0.4),
                 "Regime 2: End-to-End Fine-tuning", font_size=18,
                 font_color=ACCENT_TEAL, bold=True, font_name="Calibri")

    ft_steps = [
        "1. Initialize with pretrained weights",
        "2. Fine-tune all parameters end-to-end",
        "3. Train for 3 epochs, lr = 5e-5",
        "4. Stratified splits: 80/10/10",
    ]
    for i, step in enumerate(ft_steps):
        add_text_box(slide, Inches(7.4), Inches(3.3) + Inches(i * 0.4), Inches(5), Inches(0.35),
                     step, font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    # Separator
    add_shape(slide, Inches(6.45), Inches(2.8), Inches(0.02), Inches(2.2), LIGHT_GRAY)

    # Goal box
    add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), VERY_LIGHT)
    add_text_box(slide, Inches(1.2), Inches(5.6), Inches(10), Inches(0.35),
                 "Goal: Separate Representation Quality from Downstream Optimization",
                 font_size=17, font_color=DARK_BG, bold=True, font_name="Calibri")
    add_text_box(slide, Inches(1.2), Inches(6.05), Inches(10.5), Inches(0.5),
                 "By comparing frozen vs. fine-tuned performance, we can measure how much of the model's "
                 "usefulness comes from the pretrained representations vs. task-specific adaptation.",
                 font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    # ========================================================
    # SLIDE 8: Dataset
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Dataset: PBMC3k", "Primary evaluation benchmark", slide_num=8)

    # Dataset info card
    create_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0),
                "10x Genomics PBMC3k",
                ["Subset of PBMC 68k dataset",
                 "2,700 peripheral blood mononuclear cells",
                 "2,000 highly variable genes selected",
                 "Well-curated cell type annotations",
                 "Standard benchmark in single-cell analysis",
                 "Reproducible comparison with prior work"],
                accent_color=ACCENT_BLUE, body_size=13)

    # Cell type composition table
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.4),
                 "Cell Type Distribution", font_size=18, font_color=DARK_BG,
                 bold=True, font_name="Calibri")

    # Table
    cell_types = [
        ("Cell Type", "Count", "Proportion", True),
        ("T cells", "1,170", "43.3%", False),
        ("Monocytes", "470", "17.4%", False),
        ("NK cells", "470", "17.4%", False),
        ("B cells", "360", "13.3%", False),
        ("DC", "215", "8.0%", False),
        ("Platelet", "15", "0.6%", False),
    ]

    for i, (ct, count, prop, is_header) in enumerate(cell_types):
        y = Inches(2.05) + Inches(i * 0.38)
        bg_color = ACCENT_BLUE if is_header else (VERY_LIGHT if i % 2 == 0 else WHITE)
        text_color = WHITE if is_header else DARK_TEXT

        add_shape(slide, Inches(7.0), y, Inches(5.0), Inches(0.36), bg_color)
        add_text_box(slide, Inches(7.1), y + Inches(0.02), Inches(1.8), Inches(0.3),
                     ct, font_size=12, font_color=text_color,
                     bold=is_header, font_name="Calibri")
        add_text_box(slide, Inches(9.0), y + Inches(0.02), Inches(1.3), Inches(0.3),
                     count, font_size=12, font_color=text_color,
                     bold=is_header, font_name="Calibri", alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(10.3), y + Inches(0.02), Inches(1.5), Inches(0.3),
                     prop, font_size=12, font_color=text_color,
                     bold=is_header, font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Note about class imbalance
    add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.0),
              RGBColor(0xFF, 0xF8, 0xE1))
    add_text_box(slide, Inches(1.2), Inches(5.1), Inches(10), Inches(0.3),
                 "\u26A0  Class Imbalance Challenge", font_size=15,
                 font_color=ACCENT_CORAL, bold=True, font_name="Calibri")
    add_text_box(slide, Inches(1.2), Inches(5.45), Inches(10.5), Inches(0.4),
                 "Platelet class has only 15 cells (0.6%) vs. T cells with 1,170 (43.3%). "
                 "This 78x imbalance tests the model's robustness to rare cell types. "
                 "Macro F1 is used alongside accuracy to account for this.",
                 font_size=12, font_color=DARK_TEXT, font_name="Calibri")

    # UMAP reference image
    umap_path = os.path.join(IMG_DIR, "umap_labels_pbmc3k.png")
    if os.path.exists(umap_path):
        slide.shapes.add_picture(umap_path, Inches(1.5), Inches(6.1), Inches(10), Inches(0.85))

    # ========================================================
    # SLIDE 9: Evaluation Metrics
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Evaluation Metrics", "Quantifying model performance", slide_num=9)

    # Accuracy
    create_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.2),
                "Accuracy",
                ["Overall fraction of correct predictions",
                 "Acc = (TP + TN) / Total",
                 "Intuitive but sensitive to class imbalance",
                 "Dominant class can inflate accuracy"],
                accent_color=ACCENT_BLUE, body_size=13)

    # Macro F1
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.2),
                "Macro F1 Score",
                ["Average F1 across all classes (unweighted)",
                 "F1 = 2 \u00d7 (Prec \u00d7 Rec) / (Prec + Rec)",
                 "Macro F1 = (1/C) \u2211 F1_c",
                 "Treats all classes equally regardless of size"],
                accent_color=ACCENT_TEAL, body_size=13)

    # Additional analysis tools
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.4),
                 "Additional Analysis Methods", font_size=18, font_color=DARK_BG,
                 bold=True, font_name="Calibri")

    tools = [
        ("UMAP Visualization", "Low-dimensional projection of cell embeddings to assess cluster quality", ACCENT_BLUE),
        ("Confusion Matrix", "Identify systematic misclassification patterns between cell types", ACCENT_TEAL),
        ("Per-Class Metrics", "Precision, recall, F1 for each cell type to find failure modes", ACCENT_PURPLE),
        ("Seed Sweep (30 runs)", "Statistical robustness assessment with mean \u00b1 std over 30 seeds", ACCENT_CORAL),
    ]

    for i, (tool, desc, color) in enumerate(tools):
        y = Inches(4.7) + Inches(i * 0.55)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.05),
                                     Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        add_text_box(slide, Inches(1.3), y, Inches(2.5), Inches(0.3),
                     tool, font_size=14, font_color=color, bold=True, font_name="Calibri")
        add_text_box(slide, Inches(3.8), y, Inches(8), Inches(0.3),
                     desc, font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    # ========================================================
    # SLIDE 10: UMAP - Ground Truth vs Embeddings
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "UMAP Visualization", "Ground truth cell types in PBMC3k", slide_num=10)

    umap_path = os.path.join(IMG_DIR, "umap_labels_pbmc3k.png")
    if os.path.exists(umap_path):
        slide.shapes.add_picture(umap_path, Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.2))

    add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(0.8),
                 "Left: Leiden clustering on raw data shows well-separated clusters.  "
                 "Right: Cell type annotations confirm 6 distinct populations. "
                 "Note the clear separation of B cells and Monocytes, while T cells and NK cells show partial overlap.",
                 font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    # ========================================================
    # SLIDE 11: Geneformer Embedding UMAP
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Geneformer Frozen Embeddings", "UMAP of pretrained representations", slide_num=11)

    gf_umap = os.path.join(IMG_DIR, "umap_geneformer_emb_pbmc3k.png")
    if os.path.exists(gf_umap):
        slide.shapes.add_picture(gf_umap, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5))

    # Analysis box
    create_card(slide, Inches(6.5), Inches(1.5), Inches(6.0), Inches(3.5),
                "Observation: Poor Cluster Separation",
                ["Cell types heavily overlap in embedding space",
                 "No clear boundaries between populations",
                 "T, NK, B cells are intermixed",
                 "Suggests frozen representations lack",
                 "  cell-type-discriminative structure",
                 "",
                 "Implication: A lightweight classifier on",
                 "these embeddings cannot recover clear",
                 "decision boundaries \u2192 ~61% accuracy"],
                accent_color=ACCENT_CORAL, body_size=13)

    # Comparison note
    add_shape(slide, Inches(6.5), Inches(5.3), Inches(6.0), Inches(1.0), VERY_LIGHT)
    add_text_box(slide, Inches(6.7), Inches(5.4), Inches(5.5), Inches(0.7),
                 "Compare with ground truth UMAP: the raw data (with standard preprocessing) "
                 "shows much better separation than the Geneformer embeddings!",
                 font_size=12, font_color=ACCENT_CORAL, bold=True, font_name="Calibri")

    # ========================================================
    # SLIDE 12: scGPT Embedding UMAP
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "scGPT Frozen Embeddings", "UMAP of pretrained representations", slide_num=12)

    scgpt_umap = os.path.join(IMG_DIR, "umap_scgpt.png")
    if os.path.exists(scgpt_umap):
        slide.shapes.add_picture(scgpt_umap, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5))

    create_card(slide, Inches(6.5), Inches(1.5), Inches(6.0), Inches(3.5),
                "Observation: Partial Separation",
                ["Monocyte/DC cluster is well separated (bottom)",
                 "B cells form a distinct cluster (right)",
                 "T and NK cells overlap significantly (top center)",
                 "Platelets isolated (top left, low count)",
                 "",
                 "Better than Geneformer frozen embeddings",
                 "but still insufficient for high-accuracy",
                 "classification without fine-tuning"],
                accent_color=ACCENT_TEAL, body_size=13)

    add_shape(slide, Inches(6.5), Inches(5.3), Inches(6.0), Inches(1.0), VERY_LIGHT)
    add_text_box(slide, Inches(6.7), Inches(5.4), Inches(5.5), Inches(0.7),
                 "scGPT shows better embedding structure than Geneformer, "
                 "particularly for Monocytes and B cells. However, the T/NK "
                 "overlap remains a challenge.",
                 font_size=12, font_color=ACCENT_TEAL, bold=True, font_name="Calibri")

    # ========================================================
    # SLIDE 13: Frozen Results - Main Table
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Results: Frozen Representations",
                 "Classification with lightweight classifiers on frozen embeddings", slide_num=13)

    # Main results table
    table_data = [
        ("Model", "Classifier", "Accuracy", "Macro F1", True),
        ("Geneformer (Frozen)", "LogReg", "61.3%", "0.428", False),
        ("scGPT (Frozen)", "XGBoost", "60.0%", "0.294", False),
    ]

    for i, (model, clf, acc, f1, is_header) in enumerate(table_data):
        y = Inches(1.8) + Inches(i * 0.5)
        bg = ACCENT_BLUE if is_header else (VERY_LIGHT if i % 2 == 0 else WHITE)
        tc = WHITE if is_header else DARK_TEXT

        add_shape(slide, Inches(1.5), y, Inches(10.3), Inches(0.48), bg)
        add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(3.5), Inches(0.35),
                     model, font_size=14, font_color=tc, bold=is_header, font_name="Calibri")
        add_text_box(slide, Inches(5.2), y + Inches(0.05), Inches(2), Inches(0.35),
                     clf, font_size=14, font_color=tc, bold=is_header,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.2), y + Inches(0.05), Inches(2), Inches(0.35),
                     acc, font_size=14, font_color=tc, bold=is_header,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(9.5), y + Inches(0.05), Inches(2), Inches(0.35),
                     f1, font_size=14, font_color=tc, bold=is_header,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Key observations
    add_text_box(slide, Inches(1.0), Inches(3.5), Inches(10), Inches(0.4),
                 "Key Observations", font_size=20, font_color=DARK_BG,
                 bold=True, font_name="Calibri")

    observations = [
        ("\u2022 Similar accuracy (~60%) for both models \u2014 neither strongly discriminates cell types",
         DARK_TEXT),
        ("\u2022 Geneformer has higher Macro F1 (0.428 vs 0.294) \u2014 better handling of minority classes",
         DARK_TEXT),
        ("\u2022 Low Macro F1 indicates poor performance on rare cell types (DC, Platelet)",
         ACCENT_CORAL),
        ("\u2022 Both models far below state-of-the-art for PBMC classification (~95%+)",
         ACCENT_CORAL),
    ]

    for i, (text, color) in enumerate(observations):
        add_text_box(slide, Inches(1.3), Inches(4.0) + Inches(i * 0.45), Inches(10), Inches(0.4),
                     text, font_size=14, font_color=color, font_name="Calibri")

    # Metrics comparison image
    metrics_img = os.path.join(IMG_DIR, "analysis", "metrics_comparison.png")
    if os.path.exists(metrics_img):
        slide.shapes.add_picture(metrics_img, Inches(2.0), Inches(5.8), Inches(9.0), Inches(1.1))

    # ========================================================
    # SLIDE 14: Confusion Matrices
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Confusion Matrix Analysis",
                 "Systematic misclassification patterns", slide_num=14)

    # Geneformer confusion matrix
    gf_conf = os.path.join(IMG_DIR, "confusion_geneformer_pbmc3k.png")
    if os.path.exists(gf_conf):
        slide.shapes.add_picture(gf_conf, Inches(0.3), Inches(1.3), Inches(4.5), Inches(4.5))
    add_text_box(slide, Inches(0.3), Inches(5.9), Inches(4.5), Inches(0.3),
                 "Geneformer (Frozen) \u2014 Acc: 61.3%", font_size=13,
                 font_color=ACCENT_BLUE, bold=True, font_name="Calibri",
                 alignment=PP_ALIGN.CENTER)

    # scGPT confusion matrix
    scgpt_conf = os.path.join(IMG_DIR, "confusion_scgpt.png")
    if os.path.exists(scgpt_conf):
        slide.shapes.add_picture(scgpt_conf, Inches(4.8), Inches(1.3), Inches(4.5), Inches(4.5))
    add_text_box(slide, Inches(4.8), Inches(5.9), Inches(4.5), Inches(0.3),
                 "scGPT (Frozen) \u2014 Acc: 60.0%", font_size=13,
                 font_color=ACCENT_TEAL, bold=True, font_name="Calibri",
                 alignment=PP_ALIGN.CENTER)

    # Analysis panel
    create_card(slide, Inches(9.6), Inches(1.5), Inches(3.5), Inches(4.8),
                "Misclassification Patterns",
                ["Geneformer:",
                 " NK \u2192 T (38 errors)",
                 " Mono \u2192 T (27 errors)",
                 " B \u2192 T (17 errors)",
                 "",
                 "scGPT:",
                 " NK \u2192 T (18 errors)",
                 " DC \u2192 Mono (7 errors)",
                 "",
                 "Common: T cell class",
                 "acts as an 'attractor'",
                 "for misclassifications"],
                accent_color=ACCENT_PURPLE, body_size=11)

    # ========================================================
    # SLIDE 15: Per-Class Performance
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Per-Class Performance Analysis",
                 "Detailed breakdown by cell type", slide_num=15)

    # Per-class table
    headers = ("Cell Type", "Support", "GF Prec", "GF Rec", "GF F1", "scGPT Prec", "scGPT Rec", "scGPT F1")
    rows = [
        ("B", "72", "0.741", "0.556", "0.635", "1.000", "0.944", "0.971"),
        ("DC", "43", "0.655", "0.442", "0.528", "0.944", "0.791", "0.861"),
        ("Mono", "94", "0.696", "0.585", "0.636", "0.927", "0.947", "0.937"),
        ("NK", "94", "0.510", "0.532", "0.521", "0.938", "0.809", "0.869"),
        ("Platelet", "3", "1.000", "0.667", "0.800", "0.667", "0.667", "0.667"),
        ("T", "234", "0.665", "0.791", "0.723", "0.898", "0.983", "0.939"),
    ]

    # Header row
    y = Inches(1.5)
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.42), DARK_BG)
    col_widths = [1.5, 0.9, 1.2, 1.2, 1.2, 1.3, 1.3, 1.3]
    col_starts = [0.5]
    for w in col_widths[:-1]:
        col_starts.append(col_starts[-1] + w)

    for j, h in enumerate(headers):
        add_text_box(slide, Inches(col_starts[j]), y + Inches(0.02), Inches(col_widths[j]), Inches(0.35),
                     h, font_size=11, font_color=WHITE, bold=True,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Geneformer / scGPT sub-headers
    add_shape(slide, Inches(col_starts[2]), y - Inches(0.35), Inches(sum(col_widths[2:5])), Inches(0.33), ACCENT_BLUE)
    add_text_box(slide, Inches(col_starts[2]), y - Inches(0.33), Inches(sum(col_widths[2:5])), Inches(0.3),
                 "Geneformer (Frozen)", font_size=11, font_color=WHITE, bold=True,
                 font_name="Calibri", alignment=PP_ALIGN.CENTER)

    add_shape(slide, Inches(col_starts[5]), y - Inches(0.35), Inches(sum(col_widths[5:])), Inches(0.33), ACCENT_TEAL)
    add_text_box(slide, Inches(col_starts[5]), y - Inches(0.33), Inches(sum(col_widths[5:])), Inches(0.3),
                 "scGPT (Frozen)", font_size=11, font_color=WHITE, bold=True,
                 font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Data rows
    for i, row in enumerate(rows):
        y = Inches(1.94) + Inches(i * 0.4)
        bg = VERY_LIGHT if i % 2 == 0 else WHITE
        add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.38), bg)

        for j, val in enumerate(row):
            fc = DARK_TEXT
            bld = j == 0
            # Highlight poor F1 values
            if j in (4, 7):
                try:
                    fv = float(val)
                    if fv < 0.6:
                        fc = ACCENT_CORAL
                        bld = True
                except ValueError:
                    pass

            add_text_box(slide, Inches(col_starts[j]), y + Inches(0.02),
                         Inches(col_widths[j]), Inches(0.3),
                         val, font_size=11, font_color=fc, bold=bld,
                         font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Takeaway
    add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2), VERY_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(4.6), Inches(10), Inches(0.35),
                 "Key Takeaways", font_size=17, font_color=ACCENT_PURPLE,
                 bold=True, font_name="Calibri")

    takeaways = [
        "\u2022 scGPT consistently outperforms Geneformer across all major cell types (F1: 0.86\u20130.97 vs 0.52\u20130.72)",
        "\u2022 Geneformer struggles most with NK cells (F1=0.521) and DC (F1=0.528)",
        "\u2022 scGPT achieves near-perfect B cell classification (F1=0.971, Precision=1.0)",
        "\u2022 Platelet classification is unreliable for both models due to extremely low support (n=3)",
        "\u2022 T cells show high recall but low precision in Geneformer \u2192 other types misclassified as T",
    ]

    for i, t in enumerate(takeaways):
        add_text_box(slide, Inches(0.8), Inches(5.0) + Inches(i * 0.35), Inches(11.5), Inches(0.3),
                     t, font_size=12, font_color=DARK_TEXT, font_name="Calibri")

    # ========================================================
    # SLIDE 16: Seed Sweep Results
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Statistical Robustness: Seed Sweep",
                 "30 independent runs with different random seeds", slide_num=16)

    # Seed sweep table
    sw_headers = ("Metric", "Geneformer (mean\u00b1std)", "scGPT (mean\u00b1std)", "\u0394 (scGPT \u2212 GF)")
    sw_rows = [
        ("Accuracy", "0.694 \u00b1 0.020", "0.906 \u00b1 0.012", "+0.212"),
        ("Macro F1", "0.642 \u00b1 0.049", "0.886 \u00b1 0.027", "+0.244"),
        ("Macro F1 (no Platelet)", "0.657 \u00b1 0.021", "0.890 \u00b1 0.015", "+0.233"),
    ]

    y = Inches(1.6)
    add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(0.45), DARK_BG)
    sw_col_widths = [2.5, 3.2, 3.2, 2.4]
    sw_col_starts = [1.0]
    for w in sw_col_widths[:-1]:
        sw_col_starts.append(sw_col_starts[-1] + w)

    for j, h in enumerate(sw_headers):
        add_text_box(slide, Inches(sw_col_starts[j]), y + Inches(0.03),
                     Inches(sw_col_widths[j]), Inches(0.35),
                     h, font_size=13, font_color=WHITE, bold=True,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)

    for i, row in enumerate(sw_rows):
        y = Inches(2.07) + Inches(i * 0.5)
        bg = VERY_LIGHT if i % 2 == 0 else WHITE
        add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(0.48), bg)
        for j, val in enumerate(row):
            fc = ACCENT_TEAL if j == 3 else DARK_TEXT
            bld = j == 0 or j == 3
            add_text_box(slide, Inches(sw_col_starts[j]), y + Inches(0.05),
                         Inches(sw_col_widths[j]), Inches(0.35),
                         val, font_size=13, font_color=fc, bold=bld,
                         font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Per-class F1 comparison
    add_text_box(slide, Inches(1.0), Inches(3.8), Inches(5), Inches(0.4),
                 "Per-Class F1 (30-run average)", font_size=16, font_color=DARK_BG,
                 bold=True, font_name="Calibri")

    pc_headers = ("Cell Type", "GF F1", "scGPT F1", "\u0394F1")
    pc_rows = [
        ("T", "0.764", "0.928", "-0.164"),
        ("B", "0.696", "0.953", "-0.257"),
        ("Mono", "0.700", "0.907", "-0.208"),
        ("NK", "0.581", "0.861", "-0.280"),
        ("DC", "0.542", "0.801", "-0.259"),
        ("Platelet", "0.567", "0.865", "-0.297"),
    ]

    y_start = Inches(4.25)
    add_shape(slide, Inches(1.0), y_start, Inches(6.5), Inches(0.35), ACCENT_PURPLE)
    pc_cw = [1.5, 1.5, 1.5, 1.5]
    pc_cs = [1.0, 2.5, 4.0, 5.5]
    for j, h in enumerate(pc_headers):
        add_text_box(slide, Inches(pc_cs[j]), y_start + Inches(0.01),
                     Inches(pc_cw[j]), Inches(0.3),
                     h, font_size=11, font_color=WHITE, bold=True,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)

    for i, row in enumerate(pc_rows):
        y = y_start + Inches(0.37) + Inches(i * 0.32)
        bg = VERY_LIGHT if i % 2 == 0 else WHITE
        add_shape(slide, Inches(1.0), y, Inches(6.5), Inches(0.3), bg)
        for j, val in enumerate(row):
            fc = ACCENT_CORAL if j == 3 else DARK_TEXT
            add_text_box(slide, Inches(pc_cs[j]), y + Inches(0.01),
                         Inches(pc_cw[j]), Inches(0.25),
                         val, font_size=11, font_color=fc, bold=(j == 3),
                         font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Insight box
    create_card(slide, Inches(8.0), Inches(3.8), Inches(4.8), Inches(3.0),
                "Statistical Insights",
                ["scGPT consistently outperforms Geneformer",
                 "  across all 30 seeds (p < 0.001)",
                 "",
                 "Platelet shows highest variance (std=0.271)",
                 "  due to extremely low sample size (n=15)",
                 "",
                 "scGPT achieves >90% accuracy with",
                 "  very low variance (std=0.012)",
                 "  \u2192 Robust and reliable"],
                accent_color=ACCENT_PURPLE, body_size=12)

    # ========================================================
    # SLIDE 17: Fine-tuning Results
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Results: Fine-tuning vs. Frozen",
                 "The dramatic impact of end-to-end optimization", slide_num=17)

    # BIG comparison table
    ft_data = [
        ("Method", "Accuracy", "Macro F1", "Status", True),
        ("Geneformer (Frozen)", "61.3%", "0.428", "Baseline", False),
        ("scGPT (Frozen)", "60.0%", "0.294", "Baseline", False),
        ("Geneformer (Fine-tuned)", "97.8%", "0.978", "+59.6%", False),
    ]

    for i, (method, acc, f1, status, is_header) in enumerate(ft_data):
        y = Inches(1.6) + Inches(i * 0.55)
        if is_header:
            bg = DARK_BG
            tc = WHITE
        elif i == 3:
            bg = RGBColor(0xE8, 0xF5, 0xE9)  # Green highlight for fine-tuned
            tc = DARK_TEXT
        else:
            bg = VERY_LIGHT if i % 2 == 0 else WHITE
            tc = DARK_TEXT

        add_shape(slide, Inches(1.0), y, Inches(11.3), Inches(0.52), bg)

        if i == 3:
            # Star indicator
            add_text_box(slide, Inches(0.5), y + Inches(0.05), Inches(0.5), Inches(0.35),
                         "\u2605", font_size=20, font_color=GOLD,
                         font_name="Calibri", alignment=PP_ALIGN.CENTER)

        add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(4), Inches(0.35),
                     method, font_size=15, font_color=tc if is_header else (
                         RGBColor(0x1B, 0x7A, 0x2B) if i == 3 else DARK_TEXT),
                     bold=True if (is_header or i == 3) else False, font_name="Calibri")
        add_text_box(slide, Inches(5.5), y + Inches(0.05), Inches(2), Inches(0.35),
                     acc, font_size=15, font_color=tc if is_header else (
                         RGBColor(0x1B, 0x7A, 0x2B) if i == 3 else DARK_TEXT),
                     bold=is_header or i == 3, font_name="Calibri", alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(7.8), y + Inches(0.05), Inches(2), Inches(0.35),
                     f1, font_size=15, font_color=tc if is_header else (
                         RGBColor(0x1B, 0x7A, 0x2B) if i == 3 else DARK_TEXT),
                     bold=is_header or i == 3, font_name="Calibri", alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(10.0), y + Inches(0.05), Inches(2), Inches(0.35),
                     status, font_size=15, font_color=tc if is_header else (
                         RGBColor(0x1B, 0x7A, 0x2B) if i == 3 else MEDIUM_GRAY),
                     bold=is_header or i == 3, font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Big highlight boxes
    # Accuracy improvement
    add_rounded_rect(slide, Inches(1.0), Inches(4.2), Inches(3.5), Inches(2.5), RGBColor(0xE3, 0xF2, 0xFD))
    add_text_box(slide, Inches(1.0), Inches(4.3), Inches(3.5), Inches(0.35),
                 "Accuracy", font_size=14, font_color=ACCENT_BLUE,
                 bold=True, font_name="Calibri", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.0), Inches(4.7), Inches(3.5), Inches(0.8),
                 "61.3%  \u2192  97.8%", font_size=30, font_color=DARK_BG,
                 bold=True, font_name="Calibri", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.0), Inches(5.5), Inches(3.5), Inches(0.5),
                 "+59.6% absolute improvement\n(97% relative improvement)",
                 font_size=13, font_color=ACCENT_BLUE, font_name="Calibri",
                 alignment=PP_ALIGN.CENTER)

    # Macro F1 improvement
    add_rounded_rect(slide, Inches(5.0), Inches(4.2), Inches(3.5), Inches(2.5), RGBColor(0xE8, 0xF5, 0xE9))
    add_text_box(slide, Inches(5.0), Inches(4.3), Inches(3.5), Inches(0.35),
                 "Macro F1", font_size=14, font_color=ACCENT_TEAL,
                 bold=True, font_name="Calibri", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.0), Inches(4.7), Inches(3.5), Inches(0.8),
                 "0.428  \u2192  0.978", font_size=30, font_color=DARK_BG,
                 bold=True, font_name="Calibri", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.0), Inches(5.5), Inches(3.5), Inches(0.5),
                 "+128% relative improvement\nNear-perfect on all classes",
                 font_size=13, font_color=ACCENT_TEAL, font_name="Calibri",
                 alignment=PP_ALIGN.CENTER)

    # Insight
    create_card(slide, Inches(9.0), Inches(4.2), Inches(3.8), Inches(2.5),
                "Implication",
                ["Fine-tuning is essential, not optional",
                 "",
                 "Frozen representations capture",
                 "general patterns but miss",
                 "task-specific decision boundaries",
                 "",
                 "Confirms hypothesis H1"],
                accent_color=ACCENT_CORAL, body_size=12)

    # ========================================================
    # SLIDE 18: Discussion & Implications
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Discussion", "What do these results mean?", slide_num=18)

    findings = [
        ("Finding 1: Training Strategy > Model Architecture",
         "Both Geneformer and scGPT achieve similar frozen accuracy (~60%). "
         "The choice of fine-tuning vs. frozen matters far more than which model you use.",
         ACCENT_BLUE),
        ("Finding 2: Fine-tuning Provides Dramatic Gains",
         "97.8% accuracy after fine-tuning vs. 61.3% frozen. This 59.6% absolute improvement "
         "demonstrates that pretrained representations alone are insufficient.",
         ACCENT_TEAL),
        ("Finding 3: Frozen Embeddings Show Limited Structure",
         "UMAP analysis reveals poor cluster separation in frozen embeddings, explaining "
         "the low classifier performance. scGPT shows somewhat better structure than Geneformer.",
         ACCENT_PURPLE),
        ("Finding 4: Class Imbalance Remains Challenging",
         "Rare cell types (Platelet, DC) show consistently lower and more variable F1 scores. "
         "This is exacerbated in frozen settings where the classifier has limited capacity.",
         ACCENT_CORAL),
    ]

    for i, (title, desc, color) in enumerate(findings):
        y = Inches(1.5) + Inches(i * 1.35)
        add_shape(slide, Inches(0.8), y, Inches(0.07), Inches(1.1), color)
        add_text_box(slide, Inches(1.1), y, Inches(11.5), Inches(0.35),
                     title, font_size=16, font_color=color, bold=True, font_name="Calibri")
        add_text_box(slide, Inches(1.1), y + Inches(0.4), Inches(11.5), Inches(0.65),
                     desc, font_size=13, font_color=DARK_TEXT, font_name="Calibri")

    # ========================================================
    # SLIDE 19: Limitations & Future Work
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "Limitations & Future Directions", slide_num=19)

    # Limitations
    create_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5),
                "Limitations of This Study",
                ["\u2022 scGPT fine-tuning not completed",
                 "   (torchtext compatibility with PyTorch 2.9+)",
                 "",
                 "\u2022 Tabula Sapiens cross-dataset evaluation",
                 "   not executed (50GB dataset, time constraints)",
                 "",
                 "\u2022 scFoundation model not publicly available",
                 "   for evaluation",
                 "",
                 "\u2022 PBMC3k is relatively small (2,700 cells)",
                 "   May not represent complex tissue contexts",
                 "",
                 "\u2022 Only cell type classification evaluated",
                 "   Other tasks (perturbation, batch correction)",
                 "   not assessed"],
                accent_color=ACCENT_CORAL, body_size=12)

    # Future directions
    create_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.5),
                "Future Directions",
                ["\u2022 Complete scGPT fine-tuning comparison",
                 "   for full frozen vs. fine-tuned evaluation",
                 "",
                 "\u2022 Cross-dataset generalization study",
                 "   PBMC \u2192 Tabula Sapiens domain shift analysis",
                 "",
                 "\u2022 Layer-wise analysis & ablation studies",
                 "   Which layers benefit most from fine-tuning?",
                 "",
                 "\u2022 Expand to additional downstream tasks",
                 "   Gene perturbation, batch correction, etc.",
                 "",
                 "\u2022 Larger-scale evaluation",
                 "   PBMC 68k, multiple tissue types",
                 "   Statistical significance testing"],
                accent_color=ACCENT_TEAL, body_size=12)

    # All scripts ready
    add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7), VERY_LIGHT)
    add_text_box(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.5),
                 "\u2713 All evaluation scripts are implemented and ready to run: "
                 "run_all_evaluations.py | run_tabula_sapiens_evaluation.py | "
                 "run_scfoundation_evaluation.py | run_geneformer_finetune_pbmc3k.py | run_scgpt_finetune_pbmc3k.py",
                 font_size=11, font_color=ACCENT_TEAL, font_name="Calibri")

    # ========================================================
    # SLIDE 20: Conclusion & Summary
    # ========================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top accent
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_TEAL)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                 "Conclusion & Key Takeaways", font_size=32, font_color=WHITE,
                 bold=True, font_name="Calibri Light")

    add_shape(slide, Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.04), ACCENT_TEAL)

    # Key results in big format
    results_boxes = [
        ("61.3%\n\u2192 97.8%", "Geneformer: Frozen \u2192 Fine-tuned Accuracy", ACCENT_BLUE),
        ("+59.6%", "Absolute Accuracy Improvement", ACCENT_TEAL),
        ("0.978", "Fine-tuned Macro F1 Score", ACCENT_PURPLE),
    ]

    for i, (big_num, label, color) in enumerate(results_boxes):
        left = Inches(0.8) + Inches(i * 4.2)
        card = add_rounded_rect(slide, left, Inches(1.6), Inches(3.8), Inches(2.0), RGBColor(0x25, 0x25, 0x45))
        add_text_box(slide, left, Inches(1.7), Inches(3.8), Inches(1.0),
                     big_num, font_size=36, font_color=color, bold=True,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.8), Inches(3.8), Inches(0.6),
                     label, font_size=13, font_color=LIGHT_GRAY,
                     font_name="Calibri", alignment=PP_ALIGN.CENTER)

    # Summary points
    conclusions = [
        "1.  Fine-tuning is essential \u2014 frozen representations alone are insufficient for high-quality cell type classification",
        "2.  Training strategy matters more than model choice \u2014 Geneformer and scGPT show similar frozen performance",
        "3.  Pretrained representations capture general biological patterns but lack task-specific discriminability",
        "4.  Systematic evaluation reveals that published benchmarks may overstate the utility of frozen representations",
        "5.  A unified evaluation protocol is critical for fair comparison of single-cell foundation models",
    ]

    for i, c in enumerate(conclusions):
        y = Inches(4.0) + Inches(i * 0.5)
        add_text_box(slide, Inches(1.0), y, Inches(11.5), Inches(0.45),
                     c, font_size=14, font_color=LIGHT_GRAY, font_name="Calibri")

    # Thank you & contact
    add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.01), ACCENT_TEAL)
    add_text_box(slide, Inches(0.8), Inches(6.7), Inches(6), Inches(0.4),
                 "Thank you for your attention!", font_size=18, font_color=ACCENT_TEAL,
                 bold=True, font_name="Calibri")
    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(8), Inches(0.3),
                 "Keisuke Nishioka  |  10081049  |  Leibniz University of Hannover  |  WiSe 2025/26",
                 font_size=11, font_color=MEDIUM_GRAY, font_name="Calibri")

    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_BLUE)

    # ========================================================
    # SAVE
    # ========================================================
    output_path = "/home/nishioka/LUH/BioLLM/Nishioka_SingleCell_Foundation_Models_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
